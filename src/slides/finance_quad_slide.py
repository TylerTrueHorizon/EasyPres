"""Finance layout: 2x2 quadrant — bullets, line chart, bullets, table."""

from __future__ import annotations

import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

from intro_slide import DEFAULT_THEME
from table_slide import _estimate_table_height, SAFETY_FACTOR
from slide_chrome import (
    draw_top_bar,
    draw_logo_corner,
    draw_sources_footer,
    title_area_top,
    body_bottom_for_footer,
    font_family,
    title_text_color,
    populate_finance_bullets,
    add_section_heading,
    style_finance_data_table,
    chart_line_color,
    resolve_citation_urls_for_slide,
)


def create_finance_quad_slide(
    prs,
    title: str,
    top_left: dict,
    top_right: dict,
    bottom_left: dict,
    bottom_right: dict,
    sources_line: str = "",
    logo_text: str | None = None,
    citation_urls: dict | None = None,
    theme=None,
):
    """
    Each quadrant dict:
      top_left / bottom_left: {"heading", "bullets"}
      top_right: {"chart_title", "categories", "values"}
      bottom_right: {"heading", "headers", "rows"}
    """
    if theme is None:
        theme = DEFAULT_THEME

    cite_map = resolve_citation_urls_for_slide(
        citation_urls,
        top_left.get("bullets") or [],
        bottom_left.get("bullets") or [],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_top_bar(slide, prs, theme)
    draw_logo_corner(slide, prs, theme, logo_text)

    bottom_r = body_bottom_for_footer(theme)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    margin = int(Inches(0.45))
    gutter = int(Inches(0.2))

    t_top = title_area_top(theme)
    title_box = slide.shapes.add_textbox(
        margin,
        t_top,
        slide_w - int(Inches(1.2)),
        int(Inches(0.68)),
    )
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.text = title
    tp = tf_t.paragraphs[0]
    tp.font.bold = True
    tp.font.size = Pt(19)
    tp.font.name = font_family(theme)
    tp.font.color.rgb = title_text_color(theme)

    body_top = t_top + int(Inches(0.74))
    body_h = slide_h - body_top - bottom_r
    inner_w = slide_w - 2 * margin
    inner_h = body_h
    half_w = (inner_w - gutter) // 2
    half_h = (inner_h - gutter) // 2

    x0, x1 = margin, margin + half_w + gutter
    y0, y1 = body_top, body_top + half_h + gutter

    # --- Top left ---
    h_h = int(Inches(0.28))
    ht = slide.shapes.add_textbox(x0, y0, half_w, h_h)
    add_section_heading(ht, top_left.get("heading", ""), theme, font_pt=10)
    bb = slide.shapes.add_textbox(x0, y0 + h_h, half_w, half_h - h_h - int(Inches(0.05)))
    bb.text_frame.word_wrap = True
    populate_finance_bullets(
        bb.text_frame,
        top_left.get("bullets", []),
        theme,
        font_pt=8,
        citation_urls=cite_map,
    )

    # --- Top right: line chart ---
    cr = top_right
    ct = cr.get("chart_title") or cr.get("title", "")
    if ct:
        ctb = slide.shapes.add_textbox(x1, y0, half_w, int(Inches(0.24)))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.text = ct
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cp.font.bold = True
        cp.font.size = Pt(9)
        cp.font.name = font_family(theme)
        cp.font.color.rgb = title_text_color(theme)
        chart_y_off = int(Inches(0.26))
    else:
        chart_y_off = 0
    cats = cr.get("categories") or []
    vals = cr.get("values") or []
    cdata = CategoryChartData()
    cdata.categories = cats
    cdata.add_series("Line", vals)
    ch_top = y0 + chart_y_off
    ch_h = half_h - chart_y_off - int(Inches(0.05))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        x1,
        ch_top,
        half_w,
        ch_h,
        cdata,
    )
    ch = gf.chart
    ch.has_legend = False
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.font.name = font_family(theme)
    dl.font.size = Pt(8)
    dl.label_position = XL_LABEL_POSITION.ABOVE
    ser = plot.series[0]
    ser.format.line.color.rgb = chart_line_color(theme)
    ser.format.line.width = Pt(2)

    # --- Bottom left ---
    ht2 = slide.shapes.add_textbox(x0, y1, half_w, h_h)
    add_section_heading(ht2, bottom_left.get("heading", ""), theme, font_pt=10)
    bb2 = slide.shapes.add_textbox(x0, y1 + h_h, half_w, half_h - h_h)
    bb2.text_frame.word_wrap = True
    populate_finance_bullets(
        bb2.text_frame,
        bottom_left.get("bullets", []),
        theme,
        font_pt=8,
        citation_urls=cite_map,
    )

    # --- Bottom right: table ---
    br = bottom_right
    hdr = br.get("heading", "")
    headers = br.get("headers") or []
    rows = br.get("rows") or []
    ty = y1
    if hdr:
        htb = slide.shapes.add_textbox(x1, ty, half_w, int(Inches(0.26)))
        add_section_heading(htb, hdr, theme, font_pt=10)
        ty += int(Inches(0.28))
    if headers and rows:
        n_cols = len(headers)
        n_rows = len(rows) + 1
        tw_in = half_w / 914400
        col_ws = [tw_in / n_cols] * n_cols
        avail_in = (y1 + half_h - ty) / 914400
        data_pt = 7
        hdr_pt = 8
        for cand in range(9, 6, -1):
            tot, rhs = _estimate_table_height(headers, rows, cand, cand + 1, col_ws)
            if tot * SAFETY_FACTOR <= avail_in:
                data_pt = cand
                hdr_pt = cand + 1
                row_hs = rhs
                break
        else:
            _, row_hs = _estimate_table_height(headers, rows, 7, 8, col_ws)
            data_pt = 7
            hdr_pt = 8
        th_in = sum(row_hs)
        tbl_top = ty
        gr = slide.shapes.add_table(
            n_rows, n_cols, x1, tbl_top, half_w, Inches(th_in),
        )
        tbl = gr.table
        for i, rh in enumerate(row_hs):
            tbl.rows[i].height = Inches(rh)
        style_finance_data_table(tbl, headers, rows, theme, hdr_pt, data_pt)

    if sources_line:
        draw_sources_footer(slide, prs, theme, sources_line)
    return slide


if __name__ == "__main__":
    from themes import get_theme

    CITATION_URLS = {
        20: "https://example.com/source/20",
        21: "https://example.com/source/21",
    }

    prs = Presentation()
    ft = get_theme("finance")
    create_finance_quad_slide(
        prs,
        title="Operational Pivot Drives Cash Flow and Deleveraging",
        top_left={
            "heading": "Cash Flow & Working Capital Optimization",
            "bullets": [
                {
                    "lead": "Record Cash Generation",
                    "body": "Conversion and working capital discipline improved.",
                    "cites": [20],
                },
            ],
        },
        top_right={
            "title": "Quarterly Fleet Utilization Trend",
            "categories": ["Q4 2024", "FY 2025 Avg", "Q4 2025"],
            "values": [78.9, 79.4, 83.6],
        },
        bottom_left={
            "heading": "Deleveraging & Capital Discipline",
            "bullets": [
                {
                    "lead": "Net Leverage",
                    "body": "Path to target range on improved EBITDA.",
                    "cites": [21],
                },
            ],
        },
        bottom_right={
            "heading": "Financial Highlights: FY24 vs FY25",
            "headers": ["Metric", "FY 2024", "FY 2025", "Change"],
            "rows": [
                ["Revenue", "$1.0B", "$1.1B", "+10%"],
                ["Adj. EBITDA", "$200M", "$220M", "+10%"],
            ],
        },
        sources_line="INTRM • 10 Mar 26 [20,21]",
        citation_urls=CITATION_URLS,
        theme=ft,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    output_path = os.path.join(tests_dir, "test_finance_quad_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
