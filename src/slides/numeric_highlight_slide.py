from math import ceil
import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from intro_slide import DEFAULT_THEME


def _blend_toward_white(color, opacity=0.15):
    """Blend an RGBColor toward white at the given opacity."""
    r, g, b = color
    new_r = int(r * opacity + 255 * (1 - opacity))
    new_g = int(g * opacity + 255 * (1 - opacity))
    new_b = int(b * opacity + 255 * (1 - opacity))
    return RGBColor(new_r, new_g, new_b)


def _pick_cols(n):
    """Choose a sensible column count for n cards."""
    if n <= 3:
        return n
    if n <= 8:
        return 4
    if n <= 10:
        return 5
    return min(n, 5)


def create_numeric_highlight_slide(
    prs,
    title,
    subtitle,
    cards,
    cols=None,
    theme=None,
):
    """
    Create a numeric highlight slide with a title, subtitle, and a centered
    grid of rounded-rectangle cards each showing a label and a bold value.

    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Description text below the title
        cards: List of dicts, each with "label" and "value" keys
        cols: Number of columns (auto-picked if None)
        theme: Color theme dict (uses DEFAULT_THEME if None)
    """
    if theme is None:
        theme = DEFAULT_THEME

    n = len(cards)
    if n == 0:
        raise ValueError("cards list must not be empty")

    if cols is None:
        cols = _pick_cols(n)
    rows = ceil(n / cols)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ── Title ──────────────────────────────────────────────────────────
    title_top = Inches(0.5)
    title_height = Inches(0.75)
    title_width = Inches(8)
    title_left = (slide_width - title_width) // 2

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    MAX_TITLE_CHARS = 39
    if len(title) > MAX_TITLE_CHARS:
        title = title[:MAX_TITLE_CHARS - 3].rstrip() + "..."
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Albert Sans"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme['PRIMARY_COLOR']

    # ── Subtitle ───────────────────────────────────────────────────────
    subtitle_top = title_top + title_height + Inches(0.05)
    subtitle_height = Inches(0.7)
    subtitle_width = Inches(8)
    subtitle_left = (slide_width - subtitle_width) // 2

    subtitle_box = slide.shapes.add_textbox(
        subtitle_left, subtitle_top, subtitle_width, subtitle_height,
    )
    sf = subtitle_box.text_frame
    sf.word_wrap = True
    sf.text = subtitle
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.name = "Albert Sans"
    sp.font.size = Pt(14)
    sp.font.color.rgb = theme['NEUTRAL_DARK']

    # ── Grid geometry ──────────────────────────────────────────────────
    margin_x = Inches(0.75)
    margin_bottom = Inches(0.6)
    grid_top = subtitle_top + subtitle_height + Inches(0.35)
    grid_bottom = slide_height - margin_bottom

    available_width = slide_width - 2 * margin_x
    available_height = grid_bottom - grid_top

    gap = Inches(0.2)

    card_w = (available_width - (cols - 1) * gap) / cols
    card_h = (available_height - (rows - 1) * gap) / rows

    max_card_w = Inches(2.5)
    max_card_h = Inches(1.5)
    card_w = min(card_w, max_card_w)
    card_h = min(card_h, max_card_h)

    total_grid_w = cols * card_w + (cols - 1) * gap
    total_grid_h = rows * card_h + (rows - 1) * gap

    x_offset = margin_x + (available_width - total_grid_w) // 2
    y_offset = grid_top + (available_height - total_grid_h) // 2

    # ── Cards ──────────────────────────────────────────────────────────
    card_fill = _blend_toward_white(theme['SECONDARY_COLOR'], opacity=0.15)

    for idx, card in enumerate(cards):
        r = idx // cols
        c = idx % cols

        left = x_offset + c * (card_w + gap)
        top = y_offset + r * (card_h + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(left), int(top), int(card_w), int(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = card_fill
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        tf.paragraphs[0].space_before = Pt(6)

        label_para = tf.paragraphs[0]
        label_para.text = card["label"]
        label_para.alignment = PP_ALIGN.CENTER
        label_para.font.name = "Albert Sans"
        label_para.font.size = Pt(11)
        label_para.font.color.rgb = theme['NEUTRAL_DARK']

        value_para = tf.add_paragraph()
        value_para.text = card["value"]
        value_para.alignment = PP_ALIGN.CENTER
        value_para.font.name = "Albert Sans"
        value_para.font.size = Pt(24)
        value_para.font.bold = True
        value_para.font.color.rgb = theme['PRIMARY_COLOR']

    return slide


if __name__ == "__main__":
    prs = Presentation()

    sample_cards = [
        {"label": "Cash & Securities", "value": "$132.4 B"},
        {"label": "Op. Cash Flow", "value": "$111.5 B"},
        {"label": "Share Repurchase", "value": "$90.7 B"},
        {"label": "Term Debt", "value": "$90.7 B"},
        {"label": "R&D Exp", "value": "$34.6 B"},
        {"label": "Cash & Securities", "value": "$132.4 B"}

    ]

    create_numeric_highlight_slide(
        prs,
        title="Financial Highlights \u2013 FY2025",
        subtitle=(
            "Overview of FY2025 financial performance including strong "
            "operational cash flow, significant capital returns to shareholders "
            "via buybacks and dividends, and continued investment in R&D."
        ),
        cards=sample_cards,
        cols=4,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)

    output_path = os.path.join(tests_dir, "test_numeric_highlight_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
