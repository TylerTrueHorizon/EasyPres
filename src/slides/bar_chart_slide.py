import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from intro_slide import DEFAULT_THEME
from numeric_highlight_slide import _blend_toward_white


def _bar_color_spectrum(base_color, count):
    """Generate a list of RGBColors from darker to lighter for bar coloring."""
    if count == 1:
        return [_blend_toward_white(base_color, opacity=0.7)]
    opacities = [0.9 - i * (0.6 / (count - 1)) for i in range(count)]
    return [_blend_toward_white(base_color, opacity=op) for op in opacities]


def create_bar_chart_slide(
    prs,
    title,
    data,
    horizontal=True,
    descriptor=None,
    theme=None,
):
    """
    Create a bar chart slide with a centered title, optional descriptor, and
    a native PowerPoint bar chart (interactive during presentation).

    Args:
        prs: Presentation object
        title: Chart title text
        data: Either a dict {"Category": value, ...} for single series, or a
              list of dicts [{"name": "Series", "values": {...}}, ...] for
              multi-series
        horizontal: True for horizontal bars, False for vertical columns
        descriptor: Optional description text below the title
        theme: Color theme dict (uses DEFAULT_THEME if None)
    """
    if theme is None:
        theme = DEFAULT_THEME

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Title ──────────────────────────────────────────────────────────
    title_top = Inches(0.4)
    title_height = Inches(0.7)
    title_width = Inches(8)
    title_left = (slide_width - title_width) // 2

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    MAX_TITLE_CHARS = 39
    if len(title) > MAX_TITLE_CHARS:
        title = title[:MAX_TITLE_CHARS - 3].rstrip() + "..."
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Albert Sans"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme['PRIMARY_COLOR']

    content_top = title_top + title_height + Inches(0.05)

    # ── Descriptor (optional) ──────────────────────────────────────────
    if descriptor:
        desc_height = Inches(0.5)
        desc_width = Inches(8)
        desc_left = (slide_width - desc_width) // 2

        desc_box = slide.shapes.add_textbox(desc_left, content_top, desc_width, desc_height)
        df = desc_box.text_frame
        df.word_wrap = True
        df.text = descriptor
        dp = df.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        dp.font.name = "Albert Sans"
        dp.font.size = Pt(12)
        dp.font.color.rgb = theme['NEUTRAL_DARK']
        content_top = content_top + desc_height + Inches(0.05)

    # ── Normalize data ─────────────────────────────────────────────────
    multi_series = isinstance(data, list)
    if multi_series:
        categories = list(data[0]["values"].keys())
        series_list = data
    else:
        categories = list(data.keys())
        series_list = [{"name": "Series 1", "values": data}]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s["name"], [s["values"][c] for c in categories])

    # ── Chart placement ────────────────────────────────────────────────
    chart_width = Inches(8.5)
    chart_left = (slide_width - chart_width) // 2
    content_bottom = slide_height - Inches(0.5)
    chart_height = content_bottom - content_top
    chart_top = content_top

    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED

    graphic_frame = slide.shapes.add_chart(
        chart_type, int(chart_left), int(chart_top),
        int(chart_width), int(chart_height), chart_data,
    )
    chart = graphic_frame.chart

    # ── Legend ──────────────────────────────────────────────────────────
    if multi_series:
        chart.has_legend = True
        legend = chart.legend
        legend.include_in_layout = False
        legend.font.name = "Albert Sans"
        legend.font.size = Pt(10)
        legend.font.color.rgb = theme['NEUTRAL_DARK']
    else:
        chart.has_legend = False

    # ── Plot styling ───────────────────────────────────────────────────
    plot = chart.plots[0]
    plot.gap_width = 80

    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.show_category_name = False
    data_labels.show_series_name = False
    data_labels.font.name = "Albert Sans"
    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = theme['NEUTRAL_DARK']
    if horizontal:
        data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
    else:
        data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '#,##0'

    # ── Bar colors ─────────────────────────────────────────────────────
    if multi_series:
        base_opacities = [0.85, 0.6, 0.4, 0.25, 0.15]
        for s_idx, series in enumerate(plot.series):
            op = base_opacities[s_idx % len(base_opacities)]
            color = _blend_toward_white(theme['SECONDARY_COLOR'], opacity=op)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
    else:
        series = plot.series[0]
        n = len(categories)
        colors = _bar_color_spectrum(theme['SECONDARY_COLOR'], n)
        for i in range(n):
            point = series.points[i]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i]

    # ── Axes ───────────────────────────────────────────────────────────
    category_axis = chart.category_axis
    category_axis.visible = True
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.name = "Albert Sans"
    category_axis.tick_labels.font.size = Pt(11)
    category_axis.tick_labels.font.color.rgb = theme['NEUTRAL_DARK']
    category_axis.format.line.fill.background()

    value_axis = chart.value_axis
    value_axis.visible = True
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = theme['NEUTRAL_LIGHT']
    value_axis.tick_labels.font.name = "Albert Sans"
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.color.rgb = theme['NEUTRAL_DARK']
    value_axis.tick_labels.number_format = '#,##0'
    value_axis.format.line.fill.background()

    return slide


if __name__ == "__main__":
    prs = Presentation()

    create_bar_chart_slide(
        prs,
        title="FY2025 Mix",
        data={
            "iPhone": 209586,
            "Services": 109158,
            "Mac": 33708,
            "Wearables": 35686,
            "iPad": 28023,
        },
        horizontal=False,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)

    output_path = os.path.join(tests_dir, "test_bar_chart_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
