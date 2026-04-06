"""Finance layout: two columns of narrative + optional pie on lower right."""

from __future__ import annotations

import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


def _set_legend_manual_layout(
    leg_element,
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    """Inject <c:layout><c:manLayout> on the legend XML element for precise placement."""
    for old in leg_element.findall(qn("c:layout")):
        leg_element.remove(old)
    layout_el = parse_xml(
        f'<c:layout xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">'
        f'<c:manLayout>'
        f'<c:xMode val="factor"/>'
        f'<c:yMode val="factor"/>'
        f'<c:x val="{x}"/>'
        f'<c:y val="{y}"/>'
        f'<c:w val="{w}"/>'
        f'<c:h val="{h}"/>'
        f'</c:manLayout>'
        f'</c:layout>'
    )
    leg_element.append(layout_el)

from intro_slide import DEFAULT_THEME
from bar_chart_slide import _bar_color_spectrum
from slide_chrome import (
    draw_top_bar,
    draw_logo_corner,
    draw_sources_footer,
    title_area_top,
    body_bottom_for_footer,
    font_family,
    title_text_color,
    populate_finance_bullets,
    add_section_heading,
    resolve_citation_urls_for_slide,
)


def create_finance_dual_column_pie_slide(
    prs,
    title: str,
    columns: list[dict],
    chart: dict | None = None,
    sources_line: str = "",
    logo_text: str | None = None,
    citation_urls: dict | None = None,
    theme=None,
):
    """
    Two column bullets (section headings + finance bullets), optional pie in
    lower right of the second column.

    columns: list of {"heading": str, "bullets": [str | {"lead","body","cites"}]}
    chart: optional {"title": str, "slices": {label: float}, "show_legend": bool}
    """
    if theme is None:
        theme = DEFAULT_THEME

    cite_map = resolve_citation_urls_for_slide(
        citation_urls,
        *[c.get("bullets") or [] for c in columns],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_top_bar(slide, prs, theme)
    draw_logo_corner(slide, prs, theme, logo_text)

    bottom_r = body_bottom_for_footer(theme)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    t_top = title_area_top(theme)
    title_box = slide.shapes.add_textbox(
        int(Inches(0.55)),
        t_top,
        slide_w - int(Inches(1.35)),
        int(Inches(0.72)),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.bold = True
    tp.font.size = Pt(20)
    tp.font.name = font_family(theme)
    tp.font.color.rgb = title_text_color(theme)

    body_top = t_top + int(Inches(0.78))
    body_h = slide_h - body_top - bottom_r
    margin = int(Inches(0.5))
    gutter = int(Inches(0.35))
    col_w = (slide_w - 2 * margin - gutter) // 2
    left_l = margin
    right_l = margin + col_w + gutter

    if chart and columns and len(columns) >= 2:
        right_text_h = int(body_h * 0.4)
        chart_gap = int(Inches(0.08))
        pie_h = body_h - right_text_h - chart_gap
    else:
        right_text_h = body_h
        chart_gap = 0
        pie_h = 0

    # Left column
    if len(columns) >= 1:
        c0 = columns[0]
        sec_top = body_top
        h_tb = int(Inches(0.34))
        ht = slide.shapes.add_textbox(left_l, sec_top, col_w, h_tb)
        add_section_heading(ht, c0.get("heading", ""), theme, font_pt=11)
        bullet_top = sec_top + h_tb
        bb = slide.shapes.add_textbox(left_l, bullet_top, col_w, body_h - h_tb)
        btf = bb.text_frame
        btf.word_wrap = True
        populate_finance_bullets(
            btf, c0.get("bullets", []), theme, font_pt=10, citation_urls=cite_map,
        )

    # Right column (text)
    if len(columns) >= 2:
        c1 = columns[1]
        sec_top = body_top
        h_tb = int(Inches(0.34))
        ht = slide.shapes.add_textbox(right_l, sec_top, col_w, h_tb)
        add_section_heading(ht, c1.get("heading", ""), theme, font_pt=11)
        bullet_top = sec_top + h_tb
        bb = slide.shapes.add_textbox(
            right_l, bullet_top, col_w, max(0, right_text_h - h_tb - int(Inches(0.05)))
        )
        if right_text_h > h_tb:
            btf = bb.text_frame
            btf.word_wrap = True
            populate_finance_bullets(
                btf, c1.get("bullets", []), theme, font_pt=10, citation_urls=cite_map,
            )

    if chart:
        slices = chart.get("slices") or {}
        if slices:
            titles = chart.get("title") or ""
            c_top = body_top + right_text_h + chart_gap
            if titles:
                ct = slide.shapes.add_textbox(
                    right_l, c_top - int(Inches(0.28)), col_w, int(Inches(0.26)),
                )
                ctf = ct.text_frame
                ctf.word_wrap = True
                ctf.text = titles
                cp = ctf.paragraphs[0]
                cp.alignment = PP_ALIGN.CENTER
                cp.font.bold = True
                cp.font.size = Pt(10)
                cp.font.name = font_family(theme)
                cp.font.color.rgb = title_text_color(theme)

            cats = list(slices.keys())
            vals = list(slices.values())
            chart_data = CategoryChartData()
            chart_data.categories = cats
            chart_data.add_series("Mix", vals)
            graphic = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                right_l,
                c_top,
                col_w,
                pie_h - int(Inches(0.15)),
                chart_data,
            )
            ch = graphic.chart
            ch.has_legend = chart.get("show_legend", True)
            if ch.has_legend:
                leg = ch.legend
                leg.include_in_layout = False
                leg.font.name = font_family(theme)
                leg.font.size = Pt(8)
                # Push legend to the right edge, vertically centered in the chart area
                _set_legend_manual_layout(leg._element, x=0.76, y=0.30, w=0.24, h=0.40)
            plot = ch.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_percentage = False
            dl.show_value = True
            dl.font.name = font_family(theme)
            dl.font.size = Pt(8)
            dl.label_position = XL_LABEL_POSITION.INSIDE_END
            series = plot.series[0]
            colors = _bar_color_spectrum(theme["SECONDARY_COLOR"], len(cats))
            for i, pt in enumerate(series.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = colors[i % len(colors)]

    if sources_line:
        draw_sources_footer(slide, prs, theme, sources_line)

    return slide


if __name__ == "__main__":
    from themes import get_theme

    CITATION_URLS = {
        1: "https://example.com/source/1",
        2: "https://example.com/source/2",
        3: "https://example.com/source/3",
        4: "https://example.com/source/4",
    }

    prs = Presentation()
    ft = get_theme("finance")
    create_finance_dual_column_pie_slide(
        prs,
        title="Integrated Model Positioning for Secular Growth",
        columns=[
            {
                "heading": "Differentiated Ecosystem",
                "bullets": [
                    {
                        "lead": "Vertically Integrated Platform",
                        "body": "End-to-end capability across the value chain.",
                        "cites": [1],
                    },
                    {
                        "lead": "High Customer Stickiness",
                        "body": "Recurring relationships and renewal depth.",
                        "cites": [2, 3],
                    },
                ],
            },
            {
                "heading": "Strategic Ownership & Positioning",
                "bullets": [
                    {
                        "lead": "Controlled Company Status",
                        "body": "Aligned governance and long-term horizon.",
                        "cites": [4],
                    },
                    "Supporting narrative as a plain string bullet.",
                ],
            },
        ],
        chart={
            "title": "Revenue Mix by End-Market",
            "slices": {"Utility T&D": 60, "Infrastructure": 27, "Other": 13},
            "show_legend": True,
        },
        sources_line="10-K • 10 Mar 26 [1,2,3]; DEF 14A • 28 Apr 25 [4]",
        citation_urls=CITATION_URLS,
        theme=ft,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    output_path = os.path.join(tests_dir, "test_finance_dual_column_pie_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
