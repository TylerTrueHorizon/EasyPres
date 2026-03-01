import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from intro_slide import DEFAULT_THEME
from numeric_highlight_slide import _blend_toward_white


def create_split_bullet_slide(
    prs,
    title,
    subtitle,
    sections,
    theme=None,
):
    """
    Create a split-layout slide with a title and subtitle on the left,
    and a vertically stacked list of titled sections on the right,
    separated by subtle horizontal lines.

    Args:
        prs: Presentation object
        title: Large title text (left column)
        subtitle: Descriptor text below the title (left column)
        sections: List of dicts with "title" and "descriptor" keys
        theme: Color theme dict (uses DEFAULT_THEME if None)
    """
    if theme is None:
        theme = DEFAULT_THEME

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    desc_color = _blend_toward_white(theme['NEUTRAL_DARK'], opacity=0.55)

    # ── Left column (vertically centered) ─────────────────────────────
    left_margin = Inches(0.8)
    left_col_width = Inches(4.5)

    title_height = Inches(0.9)
    gap = Inches(0.15)
    subtitle_height = Inches(1.2)
    left_block_height = title_height + gap + subtitle_height
    left_block_top = (slide_height - left_block_height) // 2

    title_box = slide.shapes.add_textbox(
        left_margin, left_block_top, left_col_width, title_height,
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Albert Sans"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme['PRIMARY_COLOR']

    subtitle_top = left_block_top + title_height + gap

    subtitle_box = slide.shapes.add_textbox(
        left_margin, subtitle_top, left_col_width, subtitle_height,
    )
    sf = subtitle_box.text_frame
    sf.word_wrap = True
    sf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sf.text = subtitle
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sp.font.name = "Albert Sans"
    sp.font.size = Pt(13)
    sp.font.color.rgb = theme['NEUTRAL_DARK']

    # ── Right column geometry ──────────────────────────────────────────
    right_left = Inches(5.8)
    right_width = Inches(6.5)
    right_margin_end = Inches(0.5)
    usable_right_width = slide_width - right_left - right_margin_end

    content_top = Inches(0.7)
    content_bottom = slide_height - Inches(0.7)
    available_height = content_bottom - content_top

    n = len(sections)
    if n == 0:
        return slide

    section_height = available_height / n
    max_section_height = Inches(1.6)
    section_height = min(section_height, max_section_height)

    total_block_height = section_height * n
    block_top = content_top + (available_height - total_block_height) // 2

    line_height = Pt(0.5)
    line_color = RGBColor(200, 200, 200)

    # Separator lines between sections only (not above first / below last)
    for idx in range(1, n):
        line_top = block_top + section_height * idx
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(right_left), int(line_top),
            int(usable_right_width), int(line_height),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = line_color
        line.line.fill.background()
        line.shadow.inherit = False

    # Section content (single vertically-centered textbox per section)
    for idx, sec in enumerate(sections):
        sec_top = block_top + section_height * idx

        box = slide.shapes.add_textbox(
            right_left, sec_top, usable_right_width, int(section_height),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)

        title_para = tf.paragraphs[0]
        title_para.text = sec["title"]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "Albert Sans"
        title_para.font.size = Pt(15)
        title_para.font.bold = True
        title_para.font.color.rgb = theme['NEUTRAL_DARK']

        desc_para = tf.add_paragraph()
        desc_para.text = sec["descriptor"]
        desc_para.alignment = PP_ALIGN.LEFT
        desc_para.font.name = "Albert Sans"
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = desc_color
        desc_para.space_before = Pt(4)

    return slide


if __name__ == "__main__":
    prs = Presentation()

    create_split_bullet_slide(
        prs,
        title="Strategy & Growth Initiatives",
        subtitle=(
            "Focus on driving growth through premium hardware innovation, "
            "service expansion, and heavy R&D investment in AI. Bolstering "
            "supply chain resilience while expanding into key global markets."
        ),
        sections=[
            {
                "title": "Premium Hardware Cadence",
                "descriptor": (
                    "New launches for iPhone Pro, Mac, iPad, and Vision Pro "
                    "to sustain market leadership."
                ),
            },
            {
                "title": "Services Expansion",
                "descriptor": (
                    "Growth in advertising, App Store, cloud, Apple Pay, "
                    "and AppleCare revenue streams."
                ),
            },
            {
                "title": "R&D Investment",
                "descriptor": (
                    "10% YoY increase ($34.6B) targeting AI, on-device "
                    "intelligence, and spatial computing."
                ),
            },
            {
                "title": "Supply Chain Resilience",
                "descriptor": (
                    "Navigating tariffs and geopolitical risks to ensure "
                    "stability in production and distribution."
                ),
            },
            {
                "title": "Geographic Expansion",
                "descriptor": (
                    "Focusing on Europe and Japan while addressing demand "
                    "fluctuations in China."
                ),
            },
        ],
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)

    output_path = os.path.join(tests_dir, "test_split_bullet_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
