import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from intro_slide import DEFAULT_THEME
from numeric_highlight_slide import _blend_toward_white

import math

CELL_PADDING_INCHES = 0.1
CHAR_WIDTH_FACTOR = 0.52
LINE_HEIGHT_FACTOR = 1.4
SAFETY_FACTOR = 1.08


def _estimate_row_height(cells_text, font_pt, col_widths_inches):
    """Return estimated row height in inches for a single row.

    Computes wrapped line count per cell and returns the tallest cell's
    height (lines * line-height) plus vertical cell padding.
    """
    avg_char_w = font_pt * CHAR_WIDTH_FACTOR / 72
    line_h = font_pt * LINE_HEIGHT_FACTOR / 72
    max_cell_h = line_h

    for text, col_w in zip(cells_text, col_widths_inches):
        usable_w = col_w - 2 * CELL_PADDING_INCHES
        if usable_w <= 0 or avg_char_w <= 0:
            lines = 1
        else:
            chars_per_line = max(1, int(usable_w / avg_char_w))
            lines = max(1, math.ceil(len(str(text)) / chars_per_line))
        cell_h = lines * line_h
        if cell_h > max_cell_h:
            max_cell_h = cell_h

    return max_cell_h + 2 * CELL_PADDING_INCHES


def _estimate_table_height(headers, rows, data_pt, header_pt, col_widths_inches):
    """Return (total_height_inches, per_row_heights_inches) for the full table."""
    header_h = _estimate_row_height(headers, header_pt, col_widths_inches)
    row_heights = [header_h]
    for row in rows:
        row_heights.append(_estimate_row_height(row, data_pt, col_widths_inches))
    return sum(row_heights) * SAFETY_FACTOR, row_heights


def create_table_slide(
    prs,
    title,
    descriptor,
    headers,
    rows,
    theme=None,
):
    """
    Create a table slide with a bold title (top-left), a descriptor (top-right),
    and a styled table below.

    Args:
        prs: Presentation object
        title: Bold title text (top-left)
        descriptor: Description text (top-right)
        headers: List of column header strings
        rows: List of lists, each inner list is one row of cell values
        theme: Color theme dict (uses DEFAULT_THEME if None)
    """
    if theme is None:
        theme = DEFAULT_THEME

    slide_width = prs.slide_width
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Title (top-left) ───────────────────────────────────────────────
    title_left = Inches(0.75)
    title_top = Inches(0.4)
    title_width = Inches(4.5)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Albert Sans"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme['PRIMARY_COLOR']

    # ── Descriptor (top-right) ─────────────────────────────────────────
    desc_width = Inches(4.5)
    desc_left = slide_width - desc_width - Inches(0.75)
    desc_top = Inches(0.4)
    desc_height = Inches(0.8)

    desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
    df = desc_box.text_frame
    df.word_wrap = True
    df.text = descriptor
    dp = df.paragraphs[0]
    dp.alignment = PP_ALIGN.LEFT
    dp.font.name = "Albert Sans"
    dp.font.size = Pt(12)
    dp.font.color.rgb = theme['NEUTRAL_DARK']

    # ── Table ──────────────────────────────────────────────────────────
    slide_height = prs.slide_height
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_width_in = 8.5
    table_width = Inches(table_width_in)
    table_left = (slide_width - table_width) // 2

    content_top = Inches(1.5)
    content_bottom = slide_height - Inches(0.5)
    available_height = content_bottom - content_top
    available_height_in = available_height / 914400

    col_widths_in = [table_width_in / n_cols] * n_cols

    # ── Dynamic font sizing ───────────────────────────────────────────
    max_data_pt = 13
    min_data_pt = 8
    data_font_pt = min_data_pt
    header_font_pt = min_data_pt + 1
    per_row_heights = None

    for candidate in range(max_data_pt, min_data_pt - 1, -1):
        est_total, row_heights = _estimate_table_height(
            headers, rows, candidate, candidate + 1, col_widths_in,
        )
        if est_total <= available_height_in:
            data_font_pt = candidate
            header_font_pt = candidate + 1
            per_row_heights = row_heights
            break
    else:
        _, per_row_heights = _estimate_table_height(
            headers, rows, min_data_pt, min_data_pt + 1, col_widths_in,
        )

    table_height_in = sum(per_row_heights)
    table_height = Inches(table_height_in)
    table_top = content_top + (available_height - table_height) // 2

    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, int(table_left), int(table_top),
        int(table_width), int(table_height),
    )
    table = graphic_frame.table

    table.first_row = False
    table.first_col = False
    table.horz_banding = False
    table.vert_banding = False

    # ── Per-row heights ───────────────────────────────────────────────
    for i, rh in enumerate(per_row_heights):
        table.rows[i].height = Inches(rh)

    header_fill = _blend_toward_white(theme['SECONDARY_COLOR'], opacity=0.35)
    data_fill = _blend_toward_white(theme['SECONDARY_COLOR'], opacity=0.10)

    # ── Header row ─────────────────────────────────────────────────────
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Albert Sans"
        p.font.size = Pt(header_font_pt)
        p.font.bold = True
        p.font.color.rgb = theme['NEUTRAL_DARK']

    # ── Data rows ──────────────────────────────────────────────────────
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = data_fill
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Albert Sans"
            p.font.size = Pt(data_font_pt)
            p.font.color.rgb = theme['NEUTRAL_DARK']

    return slide


if __name__ == "__main__":
    prs = Presentation()

    create_table_slide(
        prs,
        title="Top Risks & Signals",
        descriptor="Condensed from FY2024 10-K Item 1A risk factors; indicators to monitor in 2025",
        headers=["Risk", "Impact", "Horizon", "Leading indicators"],
        rows=[
            ["BTC drawdown → liquidity stress", "High", "0–24 mo", "BTC price; MSTR premium/NAV; convert/ATM terms; credit spreads"],
            ["Premium compression vs spot BTC ETPs", "High", "0–36 mo", "ETP AUM growth; MSTR premium; borrow rates"],
            ["Financing shut → fixed-charge strain", "High", "0–24 mo", "Equity volatility; convert market conditions; preferred dividend settlement mix"],
            ["Fair-value + CAMT cash-tax risk", "Med–High", "12–48 mo", "BTC unrealized gains; CAMT guidance; projected AFSI"],
            ["Custody/cyber tail risk", "High", "0–60 mo", "Custodian incidents; SOC reports; regulatory actions vs custodians"],
            ["Dual-class governance concentration", "Medium", "Structural", "Insider actions; governance ratings; shareholder proposals"],
            ["Software transition/erosion", "Medium", "12–60 mo", "Renewals/deferred revenue; subscription growth; cloud gross margin"],
        ],
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)

    output_path = os.path.join(tests_dir, "test_table_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
