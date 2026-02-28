from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from pptx.dml.color import RGBColor

# Default TrueHorizon AI Color Palette
DEFAULT_THEME = {
    # Primary brand color - Deep Blue (for titles and key elements)
    'PRIMARY_COLOR': RGBColor(0, 51, 102),  # #003366
    # Secondary brand color - Bright Cyan (for accents and highlights)
    'SECONDARY_COLOR': RGBColor(0, 174, 239),  # #00AEEF
    # Tertiary color - Warm Orange (for call-to-action and emphasis)
    'TERTIARY_COLOR': RGBColor(255, 127, 0),  # #FF7F00
    # Neutral Dark - Charcoal (for body text and subtitles)
    'NEUTRAL_DARK': RGBColor(51, 51, 51),  # #333333
    # Neutral Light - Light Gray (for backgrounds and secondary text)
    'NEUTRAL_LIGHT': RGBColor(242, 242, 242)  # #F2F2F2
}



def create_intro_slide(prs, title, subtitle, date, theme=DEFAULT_THEME):
    """
    Create a simple intro slide with centered title, subtitle, date, and author.
    
    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle text
        date: Date text
        theme: Dictionary containing color theme (optional, uses DEFAULT_THEME if not provided)
    """
    
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Define vertical positions (centered with same spacing)
    # Get slide height
    slide_height = prs.slide_height

    # Calculate center position for title
    slide_width = prs.slide_width
    title_height = Inches(1)
    title_width = Inches(8)
    
    # Center the title both vertically and horizontally
    title_left = (slide_width - title_width) / 2
    title_top = (slide_height - title_height) / 2
    subtitle_top = title_top + title_height + Inches(0.1)
    
    # Add title
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    MAX_TITLE_CHARS = 30
    if len(title) > MAX_TITLE_CHARS:
        title = title[:MAX_TITLE_CHARS - 3].rstrip() + "..."
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.name = "Albert Sans"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = theme['PRIMARY_COLOR']  # Deep Blue for main title
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), subtitle_top, Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.name = "Albert Sans"
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.italic = True
    subtitle_para.font.color.rgb = theme['SECONDARY_COLOR']  # Bright Cyan for subtitle
    # Set transparency to 0.8 (20% opacity, as transparency is inverse of opacity)
    subtitle_box.fill.solid()
    subtitle_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    subtitle_box.fill.transparency = 1.0  # Make box fully transparent
    # Apply transparency to text color
    subtitle_para.font.color.brightness = 0.0
    # Note: python-pptx doesn't directly support text transparency, but we can adjust the color
    # For 80% transparency effect, we'll use a lighter version of the color
    r, g, b = theme['SECONDARY_COLOR']
    blend_factor = 0.8  # 20% opacity
    bg_r, bg_g, bg_b = 255, 255, 255  # Assuming white background
    new_r = int(r * blend_factor + bg_r * (1 - blend_factor))
    new_g = int(g * blend_factor + bg_g * (1 - blend_factor))
    new_b = int(b * blend_factor + bg_b * (1 - blend_factor))
    subtitle_para.font.color.rgb = RGBColor(new_r, new_g, new_b)
    # Add date as footer in bottom right
    date_box = slide.shapes.add_textbox(Inches(7.5), Inches(7), Inches(2), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = date
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.RIGHT
    date_para.font.name = "Albert Sans"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = theme['NEUTRAL_DARK']  # Charcoal for date
    
    return slide


if __name__ == "__main__":
    
    # Create a new presentation
    prs = Presentation()
    
    # Add an intro slide
    create_intro_slide(
        prs,
        title="Sample Presentation",
        subtitle="A Test Presentation",
        date="March 7th, 2026"
    )
    
    # Ensure the tests directory exists
    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    
    # Save the presentation
    output_path = os.path.join(tests_dir, "test_intro_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
