"""Theme-driven slide chrome: top bar, logo, sources footer, bottom bar."""

from __future__ import annotations

from typing import Any, Iterable, Optional

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from themes import font_family
from numeric_highlight_slide import _blend_toward_white


def section_heading_color(theme: dict):
    return theme.get("SECTION_HEADING_COLOR", theme["SECONDARY_COLOR"])


def citation_color(theme: dict):
    return theme.get("CITATION_COLOR", theme["SECONDARY_COLOR"])


def _url_is_citation_eligible(url: str) -> bool:
    u = (url or "").strip()
    return u.startswith("http://") or u.startswith("https://")


def _coerce_cite_key(k) -> int:
    if isinstance(k, int):
        return k
    if isinstance(k, str):
        return int(k.strip())
    return int(k)


def normalize_citation_urls(raw: Optional[dict]) -> dict[int, str]:
    """Build int-keyed URL map from None or dict (string keys allowed)."""
    if not raw:
        return {}
    out: dict[int, str] = {}
    for k, v in raw.items():
        key = _coerce_cite_key(k)
        url = str(v).strip() if v is not None else ""
        if not url:
            continue
        if not _url_is_citation_eligible(url):
            raise ValueError(f"citation_urls[{key}] must be an http(s) URL, got {url!r}")
        out[key] = url
    return out


def resolve_citation_urls_for_slide(
    citation_urls_raw: Optional[dict],
    *bullet_lists: Iterable[Any],
) -> dict[int, str]:
    """Normalize URLs and enforce coverage for all cites appearing in *bullet_lists*."""
    ids: set[int] = set()
    for bl in bullet_lists:
        ids |= collect_cite_ids_from_bullets(bl)
    urls = normalize_citation_urls(citation_urls_raw)
    validate_citation_coverage(ids, urls)
    return urls


def collect_cite_ids_from_bullets(bullets: Iterable[Any]) -> set[int]:
    """Gather citation ref ids from a list of str or dict bullets."""
    ids: set[int] = set()
    for b in bullets or []:
        if isinstance(b, dict):
            for c in b.get("cites") or []:
                ids.add(int(c))
    return ids


def collect_cite_ids_from_finance_left_blocks(blocks: Iterable[Any]) -> set[int]:
    """Cite ids from tear-sheet style left_blocks (section bullets only)."""
    ids: set[int] = set()
    for b in blocks or []:
        if not isinstance(b, dict):
            continue
        if b.get("type") == "section":
            ids |= collect_cite_ids_from_bullets(b.get("bullets") or [])
    return ids


def validate_citation_coverage(cite_ids: set[int], citation_urls: dict[int, str]) -> None:
    """Require every cited id to have a valid URL when any cites exist."""
    if not cite_ids:
        return
    missing = sorted(cite_ids - set(citation_urls.keys()))
    if missing:
        raise ValueError(
            "citation_urls must include every cited ref: "
            f"missing ref_id(s) {missing}"
        )
    for rid in cite_ids:
        if not _url_is_citation_eligible(citation_urls.get(rid, "")):
            raise ValueError(f"citation_urls[{rid}] must be a non-empty http(s) URL")


def title_text_color(theme: dict):
    return theme.get("TITLE_TEXT_COLOR", theme["NEUTRAL_DARK"])


def table_header_bg(theme: dict):
    return theme.get("TABLE_HEADER_BG", theme["PRIMARY_COLOR"])


def table_header_text(theme: dict):
    return theme.get("TABLE_HEADER_TEXT", RGBColor(255, 255, 255))


def table_first_col_bg(theme: dict):
    return theme.get("TABLE_FIRST_COL_BG", theme["NEUTRAL_LIGHT"])


def table_zebra(theme: dict, alt: bool):
    if alt:
        return theme.get("TABLE_ZEBRA_ALT", _blend_toward_white(theme["NEUTRAL_LIGHT"], opacity=0.5))
    return theme.get("TABLE_ZEBRA_BASE", RGBColor(255, 255, 255))


def chart_line_color(theme: dict):
    return theme.get("CHART_LINE_COLOR", theme["PRIMARY_COLOR"])


def top_bar_height(theme: dict) -> int:
    """Return top bar height in EMU, or 0 if disabled."""
    if not theme.get("SHOW_TOP_BAR"):
        return 0
    inches = float(theme.get("TOP_BAR_HEIGHT_IN", 0.14))
    return int(Inches(inches))


def bottom_bar_height(theme: dict, force: bool = False) -> int:
    if not force and not theme.get("SHOW_BOTTOM_BAR"):
        return 0
    inches = float(theme.get("BOTTOM_BAR_HEIGHT_IN", 0.14))
    return int(Inches(inches))


def draw_top_bar(slide, prs, theme: dict) -> None:
    if not theme.get("SHOW_TOP_BAR"):
        return
    h = top_bar_height(theme)
    color = theme.get("TOP_BAR_COLOR", theme["PRIMARY_COLOR"])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        int(prs.slide_width),
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def draw_bottom_bar(slide, prs, theme: dict) -> None:
    if not theme.get("SHOW_BOTTOM_BAR"):
        return
    h = bottom_bar_height(theme, force=True)
    top = int(prs.slide_height) - h
    color = theme.get("BOTTOM_BAR_COLOR", theme["PRIMARY_COLOR"])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        top,
        int(prs.slide_width),
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def draw_logo_corner(slide, prs, theme: dict, logo_text: Optional[str] = None) -> None:
    """Top-right logo: bold text + underline (theme-driven)."""
    if not theme.get("SHOW_LOGO", True):
        return
    text = logo_text if logo_text is not None else theme.get("LOGO_DEFAULT_TEXT", "AI")
    if not text:
        return
    slide_w = int(prs.slide_width)
    margin = Inches(0.5)
    box_w = Inches(0.85)
    box_h = Inches(0.55)
    box_l = slide_w - int(box_w) - int(margin)
    bar_h = top_bar_height(theme)
    box_top = int(bar_h + Inches(0.08))

    tb = slide.shapes.add_textbox(box_l, box_top, int(box_w), int(box_h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.name = font_family(theme)
    run.font.color.rgb = title_text_color(theme)

    line_color = theme.get("LOGO_UNDERLINE_COLOR", section_heading_color(theme))
    line_top = box_top + int(Inches(0.38))
    line_w = Inches(0.55)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        slide_w - int(margin) - int(line_w),
        line_top,
        int(line_w),
        int(Pt(3)),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = line_color
    line.line.fill.background()


def draw_sources_footer(
    slide,
    prs,
    theme: dict,
    sources_line: str,
    *,
    bottom_reserve: int = 0,
) -> None:
    """Thin rule + 'Sources: ...' above bottom_reserve EMU (e.g. bottom bar)."""
    if not sources_line:
        return
    slide_h = int(prs.slide_height)
    margin_x = Inches(0.55)
    footer_h = Inches(0.42)
    rule_y = slide_h - int(footer_h) - int(Inches(0.08)) - bottom_reserve

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(margin_x),
        rule_y,
        int(prs.slide_width) - 2 * int(margin_x),
        int(Pt(1)),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.get("FOOTER_RULE_COLOR", theme["NEUTRAL_LIGHT"])
    rule.line.fill.background()

    text_top = rule_y + int(Pt(6))
    tb = slide.shapes.add_textbox(
        int(margin_x),
        text_top,
        int(prs.slide_width) - 2 * int(margin_x),
        int(footer_h),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Sources: {sources_line}"
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.name = font_family(theme)
    p.font.color.rgb = theme.get("SOURCES_TEXT_COLOR", theme["NEUTRAL_DARK"])


def title_area_top(theme: dict) -> int:
    """Y position (EMU) for main slide title below top bar + padding."""
    return top_bar_height(theme) + int(Inches(0.2))


def body_bottom_for_footer(theme: dict, *, include_bottom_bar: bool = False) -> int:
    """Reserve EMU from bottom for sources + optional bottom bar."""
    reserve = int(Inches(0.55))
    if include_bottom_bar or theme.get("SHOW_BOTTOM_BAR"):
        reserve += bottom_bar_height(theme, force=True)
    return reserve


def _apply_citation_run_font(run, theme: dict, font_pt: int) -> None:
    run.font.name = font_family(theme)
    run.font.size = Pt(font_pt)
    run.font.color.rgb = citation_color(theme)


def _fill_finance_bullet_paragraph(
    p,
    bullet,
    theme: dict,
    font_pt: int,
    citation_urls: Optional[dict[int, str]] = None,
) -> None:
    p.level = 0
    p.space_after = Pt(2)
    urls = citation_urls or {}
    if isinstance(bullet, str):
        r = p.add_run()
        r.text = bullet
        r.font.name = font_family(theme)
        r.font.size = Pt(font_pt)
        r.font.color.rgb = theme["NEUTRAL_DARK"]
        return
    lead = str(bullet.get("lead", "") or "")
    body = str(bullet.get("body", "") or "")
    cites = bullet.get("cites") or []
    if lead:
        r1 = p.add_run()
        r1.text = f"{lead}: "
        r1.font.bold = True
        r1.font.name = font_family(theme)
        r1.font.size = Pt(font_pt)
        r1.font.color.rgb = theme["NEUTRAL_DARK"]
    r2 = p.add_run()
    r2.text = body
    r2.font.name = font_family(theme)
    r2.font.size = Pt(font_pt)
    r2.font.color.rgb = theme["NEUTRAL_DARK"]
    if cites:
        for j, c in enumerate(cites):
            cid = int(c)
            sp = p.add_run()
            sp.text = " "
            sp.font.name = font_family(theme)
            sp.font.size = Pt(font_pt)
            sp.font.color.rgb = theme["NEUTRAL_DARK"]
            cr = p.add_run()
            cr.text = f"[{cid}]"
            _apply_citation_run_font(cr, theme, font_pt)
            url = urls.get(cid)
            if url:
                cr.hyperlink.address = url


def populate_finance_bullets(
    tf,
    bullets,
    theme: dict,
    font_pt: int = 10,
    citation_urls: Optional[dict[int, str]] = None,
) -> None:
    """Replace *tf* content with one paragraph per bullet."""
    tf.word_wrap = True
    cites_map = citation_urls or {}
    tf.text = ""
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        _fill_finance_bullet_paragraph(p, bullet, theme, font_pt, cites_map)


def add_section_heading(tb, heading: str, theme: dict, font_pt: int = 12) -> None:
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = heading
    p.font.bold = True
    p.font.size = Pt(font_pt)
    p.font.name = font_family(theme)
    p.font.color.rgb = section_heading_color(theme)


def style_finance_data_table(
    table,
    headers: list,
    rows: list,
    theme: dict,
    header_pt: int,
    data_pt: int,
    *,
    first_col_body_shaded: bool = True,
    zebra_other_cols: bool = True,
) -> None:
    """Navy header, optional gray first body column, zebra on remaining columns."""
    hdr_bg = table_header_bg(theme)
    hdr_tx = table_header_text(theme)
    fc_bg = table_first_col_bg(theme)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(header_pt)
        p.font.name = font_family(theme)
        p.font.color.rgb = hdr_tx
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(data_pt)
            p.font.name = font_family(theme)
            p.font.color.rgb = theme["NEUTRAL_DARK"]
            cell.fill.solid()
            if ci == 0 and first_col_body_shaded:
                cell.fill.fore_color.rgb = fc_bg
            elif zebra_other_cols:
                cell.fill.fore_color.rgb = table_zebra(theme, alt=(ri % 2) == 1)
            else:
                cell.fill.fore_color.rgb = table_zebra(theme, alt=False)
