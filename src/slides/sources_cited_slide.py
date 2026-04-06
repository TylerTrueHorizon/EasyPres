import sys
import os
import math

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from intro_slide import DEFAULT_THEME
from themes import font_family

MIN_FONT_PT = 8
MAX_FONT_PT = 18
LINE_HEIGHT_FACTOR = 1.2
CHAR_WIDTH_FACTOR = 0.5


def _normalize_sources(sources):
    """Keep only sources that have both non-empty title and link."""
    normalized = []
    for source in sources:
        title = str(source.get("title", "")).strip()
        link = str(source.get("link", "")).strip()
        if title and link:
            normalized.append({"title": title, "link": link})
    return normalized


def _estimate_source_height(title_text, font_pt, usable_width_inches):
    """Estimate vertical height for one linked source title."""
    avg_char_w = font_pt * CHAR_WIDTH_FACTOR / 72
    if avg_char_w <= 0 or usable_width_inches <= 0:
        lines = 1
    else:
        chars_per_line = max(1, int(usable_width_inches / avg_char_w))
        lines = max(1, math.ceil(len(title_text) / chars_per_line))
    return lines * (font_pt * LINE_HEIGHT_FACTOR / 72) + (0.5 * font_pt / 72)


def create_sources_cited_slide(prs, sources, font_size=None, continued=False, theme=None):
    """
    Create one or more 'Sources Cited' slides.

    If sources don't fit on one slide at the calculated (or provided) font size,
    additional slides are created automatically.

    Args:
        prs: Presentation object
        sources: List of dicts with "title" and "link" keys.
        font_size: Optional int. If provided, forces this font size.
                   If None, auto-calculates best fit (min 8pt).
        continued: If True, adds "(cont.)" to the title of the *first* slide created.
                   Subsequent slides created by overflow always get "(cont.)".
        theme: Color theme dict.

    Returns:
        int: The font size used (so the caller can reuse it).
    """
    if theme is None:
        theme = DEFAULT_THEME

    linked_sources = _normalize_sources(sources)
    if not linked_sources:
        return font_size if font_size is not None else MIN_FONT_PT

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    margin_x = Inches(1.0)
    margin_top = Inches(1.5)
    margin_bottom = Inches(1.0)
    usable_width = slide_width - 2 * margin_x
    usable_height = slide_height - margin_top - margin_bottom

    usable_width_in = usable_width / 914400
    usable_height_in = usable_height / 914400

    if font_size is not None:
        selected_font_pt = font_size
    else:
        selected_font_pt = MIN_FONT_PT
        for candidate in range(MAX_FONT_PT, MIN_FONT_PT - 1, -1):
            total_h = 0
            for source in linked_sources:
                total_h += _estimate_source_height(source["title"], candidate, usable_width_in)
            if total_h <= usable_height_in:
                selected_font_pt = candidate
                break

    slides_content = []
    current_slide_sources = []
    current_h = 0

    for source in linked_sources:
        source_h = _estimate_source_height(source["title"], selected_font_pt, usable_width_in)
        if current_h + source_h > usable_height_in and current_slide_sources:
            slides_content.append(current_slide_sources)
            current_slide_sources = [source]
            current_h = source_h
        else:
            current_slide_sources.append(source)
            current_h += source_h

    if current_slide_sources:
        slides_content.append(current_slide_sources)

    title_color = theme.get("TITLE_TEXT_COLOR", theme["PRIMARY_COLOR"])

    for idx, content_chunk in enumerate(slides_content):
        is_first_of_batch = (idx == 0)
        show_continued = continued or (not is_first_of_batch)
        slide_title = "Sources Cited"
        if show_continued:
            slide_title += " (cont.)"

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = slide_title
        p = tf.paragraphs[0]
        p.font.name = font_family(theme)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.alignment = PP_ALIGN.CENTER

        body_box = slide.shapes.add_textbox(margin_x, margin_top, usable_width, usable_height)
        bf = body_box.text_frame
        bf.word_wrap = True

        for i, source in enumerate(content_chunk):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            run = para.add_run()
            run.text = source["title"]
            run.font.name = font_family(theme)
            run.font.size = Pt(selected_font_pt)
            run.font.color.rgb = theme["NEUTRAL_DARK"]
            run.hyperlink.address = source["link"]
            para.space_after = Pt(0.5 * selected_font_pt)

    return selected_font_pt


if __name__ == "__main__":
    prs = Presentation()
    test_sources = [
        {"title": f"Source {i}: Example Article", "link": f"https://example.com/article/{i}"}
        for i in range(1, 60)
    ]
    used_font = create_sources_cited_slide(prs, test_sources)
    print(f"Created slides with font size: {used_font}")

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    output_path = os.path.join(tests_dir, "test_sources_cited_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
