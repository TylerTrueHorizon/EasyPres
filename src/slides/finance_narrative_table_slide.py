"""Finance layout: left narrative sections + right comparison table."""

from __future__ import annotations

import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt

from intro_slide import DEFAULT_THEME
from table_slide import _estimate_row_height, _estimate_table_height, SAFETY_FACTOR
from slide_chrome import (
    draw_top_bar,
    draw_logo_corner,
    draw_sources_footer,
    title_area_top,
    body_bottom_for_footer,
    font_family,
    title_text_color,
    populate_finance_bullets,
    add_section_heading,
    style_finance_data_table,
    resolve_citation_urls_for_slide,
)


def create_finance_narrative_table_slide(
    prs,
    title: str,
    left_sections: list[dict],
    table: dict,
    sources_line: str = "",
    logo_text: str | None = None,
    citation_urls: dict | None = None,
    theme=None,
):
    """
    left_sections: [{"heading": str, "bullets": [...]}, ...]
    table: {"heading": str, "headers": [...], "rows": [[...]]}
    """
    if theme is None:
        theme = DEFAULT_THEME

    cite_map = resolve_citation_urls_for_slide(
        citation_urls,
        *[sec.get("bullets") or [] for sec in left_sections],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_top_bar(slide, prs, theme)
    draw_logo_corner(slide, prs, theme, logo_text)

    bottom_r = body_bottom_for_footer(theme)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    t_top = title_area_top(theme)
    title_box = slide.shapes.add_textbox(
        int(Inches(0.55)),
        t_top,
        slide_w - int(Inches(1.35)),
        int(Inches(0.72)),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.bold = True
    tp.font.size = Pt(20)
    tp.font.name = font_family(theme)
    tp.font.color.rgb = title_text_color(theme)

    body_top = t_top + int(Inches(0.78))
    body_h = slide_h - body_top - bottom_r
    margin = int(Inches(0.45))
    gutter = int(Inches(0.32))
    left_w = int(slide_w * 0.43)
    right_l = margin + left_w + gutter
    right_w = slide_w - right_l - margin

    # --- Left narrative ---
    y = body_top
    body_bottom = body_top + body_h
    for sec in left_sections:
        if y >= body_bottom - int(Inches(0.3)):
            break
        head_h = int(Inches(0.32))
        ht = slide.shapes.add_textbox(margin, y, left_w, head_h)
        ht.text_frame.word_wrap = True
        add_section_heading(ht, sec.get("heading", ""), theme, font_pt=11)
        y += head_h
        bullets = sec.get("bullets", [])
        # 0.38 in per bullet accounts for typical wrapping at 9pt in ~4-inch column
        raw_h = int(Inches(0.2) + len(bullets) * int(Inches(0.38)))
        remaining = body_bottom - y - int(Inches(0.12))
        block_h = max(int(Inches(0.3)), min(raw_h, remaining))
        bb = slide.shapes.add_textbox(margin, y, left_w, block_h)
        bb.text_frame.word_wrap = True
        populate_finance_bullets(
            bb.text_frame, bullets, theme, font_pt=9, citation_urls=cite_map,
        )
        y += block_h + int(Inches(0.1))

    # --- Right table ---
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers:
        if sources_line:
            draw_sources_footer(slide, prs, theme, sources_line)
        return slide

    th_title = table.get("heading", "")
    t_heading_top = body_top
    th_h = int(Inches(0.3)) if th_title else 0
    if th_title:
        thb = slide.shapes.add_textbox(right_l, t_heading_top, right_w, th_h)
        thb.text_frame.word_wrap = True
        add_section_heading(thb, th_title, theme, font_pt=11)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_w_in = right_w / 914400
    col_widths_in = [table_w_in / n_cols] * n_cols
    content_top = body_top + th_h + int(Inches(0.06))
    available_height = body_top + body_h - content_top
    available_height_in = available_height / 914400

    max_data_pt = 11
    min_data_pt = 7
    data_font_pt = min_data_pt
    header_font_data = min_data_pt + 1
    per_row_heights = None
    for candidate in range(max_data_pt, min_data_pt - 1, -1):
        est_total, row_heights = _estimate_table_height(
            headers, rows, candidate, candidate + 1, col_widths_in,
        )
        if est_total * SAFETY_FACTOR <= available_height_in:
            data_font_pt = candidate
            header_font_data = candidate + 1
            per_row_heights = row_heights
            break
    if per_row_heights is None:
        _, per_row_heights = _estimate_table_height(
            headers, rows, min_data_pt, min_data_pt + 1, col_widths_in,
        )

    table_height_in = sum(per_row_heights)
    table_height = Inches(table_height_in)
    table_width = Inches(table_w_in)
    table_left = right_l
    table_top = int(content_top)

    graphic = slide.shapes.add_table(
        n_rows, n_cols, table_left, table_top, int(table_width), int(table_height),
    )
    tbl = graphic.table
    for i, rh in enumerate(per_row_heights):
        tbl.rows[i].height = Inches(rh)
    style_finance_data_table(
        tbl, headers, rows, theme, header_font_data, data_font_pt,
    )

    if sources_line:
        draw_sources_footer(slide, prs, theme, sources_line)
    return slide


if __name__ == "__main__":
    from themes import get_theme

    CITATION_URLS = {
        7: "https://example.com/source/7",
        8: "https://example.com/source/8",
        9: "https://example.com/source/9",
    }

    prs = Presentation()
    ft = get_theme("finance")
    create_finance_narrative_table_slide(
        prs,
        title="Vertical Integration Drives Structural Cost Advantages",
        left_sections=[
            {
                "heading": "Manufacturing-Led Sourcing Moat",
                "bullets": [
                    {
                        "lead": "Sourcing Leverage",
                        "body": "Direct manufacturing scale and purchasing power.",
                        "cites": [7],
                    },
                    {
                        "lead": "Superior Economics",
                        "body": "Higher margins vs. rental-only peers.",
                        "cites": [8],
                    },
                ],
            },
            {
                "heading": "Product Innovation & ESG Wins",
                "bullets": [
                    {
                        "lead": "Electric PTO Technology",
                        "body": "Differentiated low-emission offerings.",
                        "cites": [9],
                    },
                ],
            },
        ],
        table={
            "heading": "CTOS Vertical Integration vs. Traditional Rental",
            "headers": ["Business Model", "Sourcing", "ROIC & Margin"],
            "rows": [
                ["Platform", "Manufacturing-led", "High-teens ROIC"],
                ["Traditional", "Third-party", "Lower returns"],
            ],
        },
        sources_line="10-K • Mar 26 [7,8,9]",
        citation_urls=CITATION_URLS,
        theme=ft,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    output_path = os.path.join(tests_dir, "test_finance_narrative_table_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
