import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from intro_slide import DEFAULT_THEME
from numeric_highlight_slide import _blend_toward_white

import math


def _estimate_wrapped_lines(text, font_pt, usable_width_inches):
    """Estimate how many wrapped lines a string occupies at a given font size."""
    avg_char_width_inches = font_pt * 0.52 / 72
    if avg_char_width_inches <= 0:
        return 1
    chars_per_line = max(1, int(usable_width_inches / avg_char_width_inches))
    return max(1, math.ceil(len(text) / chars_per_line))


def _compute_bullet_font_size(cards, card_w_inches, card_h_inches,
                               title_font_pt, margin_lr_inches, margin_tb_inches):
    """Find the largest bullet font size that fits every card without overflow.

    Returns an integer point size. Bullet size is capped at 80% of title_font_pt.
    """
    usable_w = card_w_inches - 2 * margin_lr_inches
    usable_h = (card_h_inches - 2 * margin_tb_inches) * 0.92

    max_bullet_pt = int(title_font_pt * 0.8)
    min_bullet_pt = 8

    for candidate in range(max_bullet_pt, min_bullet_pt - 1, -1):
        fits_all = True
        for card in cards:
            title_lines = _estimate_wrapped_lines(card["title"], title_font_pt, usable_w)
            title_height = title_lines * (title_font_pt * 1.3 / 72)
            spacing_after_title = 8 / 72  # Pt(8) in inches

            bullet_height = 0
            for i, bt in enumerate(card.get("bullets", [])):
                full_text = "\u2022  " + bt
                lines = _estimate_wrapped_lines(full_text, candidate, usable_w)
                bullet_height += lines * (candidate * 1.3 / 72)
                if i > 0:
                    bullet_height += 6 / 72  # Pt(6) space_before

            total = title_height + spacing_after_title + bullet_height
            if total > usable_h:
                fits_all = False
                break

        if fits_all:
            return candidate

    return min_bullet_pt


def create_bulleted_boxes_slide(
    prs,
    title,
    cards,
    theme=None,
):
    """
    Create a slide with a centered title and a horizontal row of rounded-
    rectangle cards, each containing a bold title and bullet points.
    Cards alternate between white-with-border and tinted fill.

    Args:
        prs: Presentation object
        title: Centered title text at the top
        cards: List of dicts with "title" and "bullets" keys,
               e.g. [{"title": "Risks", "bullets": ["Item 1", "Item 2"]}, ...]
        theme: Color theme dict (uses DEFAULT_THEME if None)
    """
    if theme is None:
        theme = DEFAULT_THEME

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tinted_fill = _blend_toward_white(theme['SECONDARY_COLOR'], opacity=0.15)
    border_color = RGBColor(200, 200, 200)

    # ── Scaling by card count ─────────────────────────────────────────
    n = len(cards)
    if n == 0:
        pass  # will return after drawing slide title
    if n <= 2:
        circle_size_in, number_font_pt, number_height_in, card_title_pt, slide_title_pt = 0.80, 28, 0.90, 26, 32
    elif n == 3:
        circle_size_in, number_font_pt, number_height_in, card_title_pt, slide_title_pt = 0.65, 22, 0.75, 22, 32
    elif n == 4:
        circle_size_in, number_font_pt, number_height_in, card_title_pt, slide_title_pt = 0.55, 18, 0.65, 18, 30
    else:
        circle_size_in, number_font_pt, number_height_in, card_title_pt, slide_title_pt = 0.45, 14, 0.55, 16, 28

    # ── Title ──────────────────────────────────────────────────────────
    title_top = Inches(0.5)
    title_height = Inches(0.9)
    title_width = Inches(8)
    title_left = (slide_width - title_width) // 2

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Albert Sans"
    p.font.size = Pt(slide_title_pt)
    p.font.bold = True
    p.font.color.rgb = theme['PRIMARY_COLOR']

    if n == 0:
        return slide

    # ── Card geometry ──────────────────────────────────────────────────
    margin_x = Inches(0.75)
    gap = Inches(0.3)
    number_height = Inches(number_height_in)
    content_top = title_top + title_height + Inches(0.35)
    content_bottom = slide_height - Inches(0.6)
    available_width = slide_width - 2 * margin_x
    available_height = content_bottom - content_top

    card_w = (available_width - (n - 1) * gap) / n
    card_w = min(card_w, Inches(4.5))

    card_h = available_height - number_height - Inches(0.1)
    card_h = min(card_h, Inches(4.2))

    total_row_w = n * card_w + (n - 1) * gap
    x_offset = margin_x + (available_width - total_row_w) // 2
    number_y = content_top
    card_y = number_y + number_height + Inches(0.1)

    number_color = _blend_toward_white(theme['SECONDARY_COLOR'], opacity=0.20)

    margin_lr_in = 0.15
    margin_tb_in = 0.15
    bullet_pt = _compute_bullet_font_size(
        cards,
        card_w_inches=card_w / 914400,
        card_h_inches=card_h / 914400,
        title_font_pt=card_title_pt,
        margin_lr_inches=margin_lr_in,
        margin_tb_inches=margin_tb_in,
    )

    # ── Cards ──────────────────────────────────────────────────────────
    for idx, card in enumerate(cards):
        col_left = x_offset + idx * (card_w + gap)

        # Number circle above each card
        circle_size = Inches(circle_size_in)
        circle_left = col_left + Inches(0.1)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(circle_left), int(number_y),
            int(circle_size), int(circle_size),
        )
        circle.shadow.inherit = False
        circle.fill.solid()
        circle.fill.fore_color.rgb = number_color
        circle.line.fill.background()
        ntf = circle.text_frame
        ntf.text = str(idx + 1)
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        np.font.name = "Albert Sans"
        np.font.size = Pt(number_font_pt)
        np.font.bold = True
        np.font.color.rgb = theme['NEUTRAL_DARK']

        # Card rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(col_left), int(card_y), int(card_w), int(card_h),
        )
        shape.shadow.inherit = False
        shape.adjustments[0] = 0.02

        if idx % 2 == 0:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = tinted_fill
            shape.line.fill.background()

        ctf = shape.text_frame
        ctf.word_wrap = True
        ctf.auto_size = None
        ctf.vertical_anchor = MSO_ANCHOR.TOP
        ctf.margin_left = Inches(0.15)
        ctf.margin_right = Inches(0.15)
        ctf.margin_top = Inches(0.15)
        ctf.margin_bottom = Inches(0.15)

        title_para = ctf.paragraphs[0]
        title_para.text = card["title"]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = "Albert Sans"
        title_para.font.size = Pt(card_title_pt)
        title_para.font.bold = True
        title_para.font.color.rgb = theme['NEUTRAL_DARK']
        title_para.space_after = Pt(8)

        for bullet_text in card.get("bullets", []):
            bp = ctf.add_paragraph()
            bp.text = "\u2022  " + bullet_text
            bp.alignment = PP_ALIGN.LEFT
            bp.font.name = "Albert Sans"
            bp.font.size = Pt(bullet_pt)
            bp.font.color.rgb = theme['NEUTRAL_DARK']
            bp.space_before = Pt(6)

            hang = int(Pt(bullet_pt)*.9)
            pPr = bp._p.get_or_add_pPr()
            pPr.set('marL', str(hang))
            pPr.set('indent', str(-hang))

        if n >= 5 or bullet_pt == 8:
            ctf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    return slide


if __name__ == "__main__":
    prs = Presentation()

    create_bulleted_boxes_slide(
        prs,
        title="Risks & Key Debates",
        cards=[
            {
                "title": "Regulatory Risks",
                "bullets": [
                    "EU DMA and antitrust measures may compress Services take-rates",
                    "Greater China sales decline (-4%) raises concerns",
                    "AI integration brings safety challenges",
                ],
            },
            {
                "title": "Financial Risks",
                "bullets": [
                    "New tariffs may erode Product GM (36.8%)",
                    "High share-repurchase program increases reliance on free cash flow",
                    "Currency headwinds from strong USD",
                ],
            }
        ],
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)

    output_path = os.path.join(tests_dir, "test_bulleted_boxes_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
