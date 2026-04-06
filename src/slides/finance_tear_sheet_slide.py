"""Finance layout: dense left stack + right horizontal bars + bottom navy bar."""

from __future__ import annotations

import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

from intro_slide import DEFAULT_THEME
from themes import merge_theme
from table_slide import _estimate_table_height, SAFETY_FACTOR
from bar_chart_slide import _bar_color_spectrum
from slide_chrome import (
    draw_top_bar,
    draw_bottom_bar,
    draw_logo_corner,
    draw_sources_footer,
    title_area_top,
    body_bottom_for_footer,
    bottom_bar_height,
    font_family,
    title_text_color,
    populate_finance_bullets,
    add_section_heading,
    style_finance_data_table,
    collect_cite_ids_from_finance_left_blocks,
    normalize_citation_urls,
    validate_citation_coverage,
)


def create_finance_tear_sheet_slide(
    prs,
    title: str,
    left_blocks: list[dict],
    right_chart: dict,
    sources_line: str = "",
    logo_text: str | None = None,
    citation_urls: dict | None = None,
    theme=None,
):
    """
    left_blocks: [{"type": "section", "heading", "bullets"}, {"type": "table", ...}]
    right_chart: {"title", "categories", "values"}
    Enables SHOW_BOTTOM_BAR on a merged theme for this slide only.
    """
    if theme is None:
        theme = DEFAULT_THEME
    theme = merge_theme(theme, {"SHOW_BOTTOM_BAR": True})

    cite_ids = collect_cite_ids_from_finance_left_blocks(left_blocks)
    cite_map = normalize_citation_urls(citation_urls)
    validate_citation_coverage(cite_ids, cite_map)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_top_bar(slide, prs, theme)
    draw_logo_corner(slide, prs, theme, logo_text)

    bottom_r = body_bottom_for_footer(theme, include_bottom_bar=True)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    margin = int(Inches(0.45))
    left_w = int(slide_w * 0.42)
    gutter = int(Inches(0.55))
    right_l = margin + left_w + gutter
    right_w = slide_w - right_l - int(Inches(0.3))

    t_top = title_area_top(theme)
    title_box = slide.shapes.add_textbox(
        margin,
        t_top,
        slide_w - int(Inches(1.15)),
        int(Inches(0.65)),
    )
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.text = title
    tp = tf_t.paragraphs[0]
    tp.font.bold = True
    tp.font.size = Pt(18)
    tp.font.name = font_family(theme)
    tp.font.color.rgb = title_text_color(theme)

    body_top = t_top + int(Inches(0.72))
    body_h = slide_h - body_top - bottom_r
    y = body_top
    bottom_limit = body_top + body_h

    for block in left_blocks:
        if y >= bottom_limit - int(Inches(0.3)):
            break
        btype = block.get("type", "section")
        if btype == "section":
            h_h = int(Inches(0.26))
            ht = slide.shapes.add_textbox(margin, y, left_w, h_h)
            add_section_heading(ht, block.get("heading", ""), theme, font_pt=10)
            y += h_h
            bullets = block.get("bullets", [])
            est_h = int(Inches(0.22)) * max(1, min(len(bullets), 4)) + int(Inches(0.15))
            est_h = min(est_h, bottom_limit - y - int(Inches(0.1)))
            bb = slide.shapes.add_textbox(margin, y, left_w, est_h)
            bb.text_frame.word_wrap = True
            populate_finance_bullets(
                bb.text_frame, bullets, theme, font_pt=8, citation_urls=cite_map,
            )
            y += est_h + int(Inches(0.08))
        elif btype == "table":
            hdr_t = block.get("heading", "")
            if hdr_t:
                h_h = int(Inches(0.26))
                ht = slide.shapes.add_textbox(margin, y, left_w, h_h)
                add_section_heading(ht, hdr_t, theme, font_pt=10)
                y += h_h
            headers = block.get("headers") or []
            rows = block.get("rows") or []
            if not headers or y >= bottom_limit:
                continue
            n_cols = len(headers)
            n_rows = len(rows) + 1
            tw_in = left_w / 914400
            col_ws = [tw_in / n_cols] * n_cols
            avail_in = (bottom_limit - y) / 914400
            data_pt = 6
            hdr_pt = 7
            row_hs = None
            for cand in range(8, 5, -1):
                tot, rhs = _estimate_table_height(headers, rows, cand, cand + 1, col_ws)
                if tot * SAFETY_FACTOR <= avail_in * 0.95:
                    data_pt = cand
                    hdr_pt = cand + 1
                    row_hs = rhs
                    break
            if row_hs is None:
                _, row_hs = _estimate_table_height(headers, rows, 6, 7, col_ws)
            th_in = sum(row_hs)
            gr = slide.shapes.add_table(
                n_rows, n_cols, margin, y, left_w, Inches(th_in),
            )
            tbl = gr.table
            for i, rh in enumerate(row_hs):
                tbl.rows[i].height = Inches(rh)
            style_finance_data_table(tbl, headers, rows, theme, hdr_pt, data_pt)
            y += int(Inches(th_in)) + int(Inches(0.1))

    # Right chart
    rch = right_chart or {}
    ctitle = rch.get("title", "")
    ch_top = body_top
    if ctitle:
        ctb = slide.shapes.add_textbox(
            right_l, ch_top, right_w, int(Inches(0.55)),
        )
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.text = ctitle
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cp.font.bold = True
        cp.font.size = Pt(10)
        cp.font.name = font_family(theme)
        cp.font.color.rgb = title_text_color(theme)
        ch_top += int(Inches(0.52))

    cats = rch.get("categories") or []
    vals = rch.get("values") or []
    c_h = bottom_limit - ch_top - int(Inches(0.1))
    cdata = CategoryChartData()
    cdata.categories = cats
    cdata.add_series("Vals", vals)
    chart_indent = int(Inches(0.35))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        right_l + chart_indent,
        ch_top,
        right_w - chart_indent,
        c_h,
        cdata,
    )
    ch = gf.chart
    ch.has_legend = False
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.font.name = font_family(theme)
    dl.font.size = Pt(8)
    dl.font.color.rgb = theme.get("NEUTRAL_LIGHT", theme["NEUTRAL_DARK"])
    # INSIDE_END keeps labels inside bars — prevents negative-bar labels
    # from leaking left out of the chart boundary into the text column.
    dl.label_position = XL_LABEL_POSITION.INSIDE_END
    series = plot.series[0]
    n = len(cats)
    colors = _bar_color_spectrum(theme["PRIMARY_COLOR"], max(1, n))
    for i in range(n):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]

    cat_ax = ch.category_axis
    cat_ax.tick_labels.font.name = font_family(theme)
    cat_ax.tick_labels.font.size = Pt(8)
    cat_ax.tick_labels.font.color.rgb = theme["NEUTRAL_DARK"]

    val_ax = ch.value_axis
    val_ax.tick_labels.font.name = font_family(theme)
    val_ax.tick_labels.font.size = Pt(8)
    val_ax.tick_labels.font.color.rgb = theme["NEUTRAL_DARK"]
    val_ax.major_gridlines.format.line.color.rgb = theme["NEUTRAL_LIGHT"]

    if sources_line:
        draw_sources_footer(
            slide, prs, theme, sources_line,
            bottom_reserve=bottom_bar_height(theme, force=True),
        )

    draw_bottom_bar(slide, prs, theme)
    return slide


if __name__ == "__main__":
    from themes import get_theme

    CITATION_URLS = {
        1: "https://example.com/10k-mar-26",
        2: "https://example.com/deutsche-bank-initiation",
        3: "https://example.com/oppenheimer-upgrade",
        4: "https://example.com/ms-equity-note",
        5: "https://example.com/barclays-comp-analysis",
        6: "https://example.com/bloomberg-grid-survey",
        7: "https://example.com/eia-power-outlook",
        8: "https://example.com/company-ir-day",
    }

    prs = Presentation()
    ft = get_theme("finance")
    create_finance_tear_sheet_slide(
        prs,
        title="AAPL: Valuation, Catalysts & Downside Risks — Comprehensive Tear Sheet",
        left_blocks=[
            {
                "type": "section",
                "heading": "Bull Case: Structural Tailwinds",
                "bullets": [
                    {
                        "lead": "Grid Modernization Supercycle",
                        "body": "Decade-long transmission replacement cycle accelerating; AAPL uniquely positioned with vertically integrated power management silicon.",
                        "cites": [1, 6],
                    },
                    {
                        "lead": "AI & Data Center Demand",
                        "body": "Hyperscaler capex growing 30%+ YoY; on-device inference drives upgrade urgency across installed base of 2.2B active devices.",
                        "cites": [4],
                    },
                    {
                        "lead": "Services Flywheel",
                        "body": "High-margin services (~74% gross margin) growing 15%+ and compounding; App Store, iCloud, and Apple Intelligence subscriptions driving durable ARPU expansion.",
                        "cites": [1],
                    },
                ],
            },
            {
                "type": "section",
                "heading": "Bear Case: Key Risks to Monitor",
                "bullets": [
                    {
                        "lead": "China Revenue Concentration",
                        "body": "~19% of revenue exposed to geopolitical risk; Huawei share recovery and regulatory friction remain overhangs.",
                        "cites": [2],
                    },
                    {
                        "lead": "Valuation Premium Compression",
                        "body": "Trading at 28x forward P/E vs. Mag-7 median of 24x; multiple at risk if AI monetization timeline slips beyond FY26.",
                        "cites": [3, 5],
                    },
                    {
                        "lead": "Regulatory Headwinds",
                        "body": "DOJ antitrust action on App Store economics and EU DMA compliance could reduce Services segment take-rate by 2–4pp.",
                        "cites": [2],
                    },
                ],
            },
            {
                "type": "section",
                "heading": "Near-Term Catalysts",
                "bullets": [
                    {
                        "lead": "iPhone 17 Launch (Sep 26)",
                        "body": "Foldable form factor and Apple Intelligence Pro tier expected to drive largest upgrade cycle since iPhone 12.",
                        "cites": [8],
                    },
                    {
                        "lead": "Capital Return Program",
                        "body": "$110B buyback authorization; dividend yield at 0.5% with potential for special dividend if net cash exceeds $100B.",
                        "cites": [1],
                    },
                ],
            },
            {
                "type": "table",
                "heading": "Street Price Target Summary",
                "headers": ["Firm", "Rating", "Method", "Multiple", "Target"],
                "rows": [
                    ["Deutsche Bank", "Buy", "DCF + Comps", "29x FY26E", "$245"],
                    ["Oppenheimer", "Outperform", "EV/EBITDA", "18x FY26E", "$240"],
                    ["Morgan Stanley", "Overweight", "SOTP", "—", "$255"],
                    ["Barclays", "Equal Weight", "P/E Comps", "26x FY26E", "$210"],
                    ["Goldman Sachs", "Buy", "DCF", "28x FY26E", "$238"],
                ],
            },
        ],
        right_chart={
            "title": "12M Price Target vs. Current & Peers",
            "categories": [
                "GS target",
                "DB target",
                "MS target",
                "Consensus",
                "Current price",
                "Barclays target",
                "52-wk low",
            ],
            "values": [238, 245, 255, 232, 198, 210, 164],
        },
        sources_line=(
            "10-K Mar 26 [1]; Deutsche Bank [2]; Oppenheimer [3]; "
            "Morgan Stanley [4]; Barclays [5]; Bloomberg Grid Survey [6]; "
            "EIA Power Outlook [7]; AAPL IR Day [8]"
        ),
        citation_urls=CITATION_URLS,
        theme=ft,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    output_path = os.path.join(tests_dir, "test_finance_tear_sheet_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
