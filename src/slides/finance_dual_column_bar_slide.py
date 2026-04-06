"""Finance layout: two columns of bullets above a full-width column chart."""

from __future__ import annotations

import sys
import os

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

from intro_slide import DEFAULT_THEME
from bar_chart_slide import _bar_color_spectrum
from slide_chrome import (
    draw_top_bar,
    draw_logo_corner,
    draw_sources_footer,
    title_area_top,
    body_bottom_for_footer,
    font_family,
    title_text_color,
    populate_finance_bullets,
    add_section_heading,
    resolve_citation_urls_for_slide,
)


def create_finance_dual_column_bar_slide(
    prs,
    title: str,
    columns: list[dict],
    chart: dict,
    sources_line: str = "",
    logo_text: str | None = None,
    citation_urls: dict | None = None,
    theme=None,
):
    """
    columns: [{"heading", "bullets"}, ...] typically length 2
    chart: {"title", "categories": [...], "values": [...]}
    """
    if theme is None:
        theme = DEFAULT_THEME

    cite_map = resolve_citation_urls_for_slide(
        citation_urls,
        *[c.get("bullets") or [] for c in columns],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_top_bar(slide, prs, theme)
    draw_logo_corner(slide, prs, theme, logo_text)

    bottom_r = body_bottom_for_footer(theme)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    t_top = title_area_top(theme)
    title_box = slide.shapes.add_textbox(
        int(Inches(0.55)),
        t_top,
        slide_w - int(Inches(1.35)),
        int(Inches(0.7)),
    )
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.text = title
    tp = tf_t.paragraphs[0]
    tp.font.bold = True
    tp.font.size = Pt(20)
    tp.font.name = font_family(theme)
    tp.font.color.rgb = title_text_color(theme)

    body_top = t_top + int(Inches(0.76))
    text_band_h = int(slide_h * 0.34)
    chart_top = body_top + text_band_h + int(Inches(0.15))
    chart_h = slide_h - chart_top - bottom_r - int(Inches(0.2))

    margin = int(Inches(0.5))
    gutter = int(Inches(0.35))
    col_w = (slide_w - 2 * margin - gutter) // 2
    for idx, col in enumerate(columns[:2]):
        x = margin if idx == 0 else margin + col_w + gutter
        h_tb = int(Inches(0.3))
        ht = slide.shapes.add_textbox(x, body_top, col_w, h_tb)
        add_section_heading(ht, col.get("heading", ""), theme, font_pt=11)
        bb = slide.shapes.add_textbox(x, body_top + h_tb, col_w, text_band_h - h_tb - int(Inches(0.05)))
        bb.text_frame.word_wrap = True
        populate_finance_bullets(
            bb.text_frame,
            col.get("bullets", []),
            theme,
            font_pt=9,
            citation_urls=cite_map,
        )

    cats = chart.get("categories") or []
    vals = chart.get("values") or []
    chart_title = chart.get("title") or ""
    if chart_title:
        ctb = slide.shapes.add_textbox(
            margin,
            chart_top - int(Inches(0.28)),
            slide_w - 2 * margin,
            int(Inches(0.26)),
        )
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.text = chart_title
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cp.font.bold = True
        cp.font.size = Pt(11)
        cp.font.name = font_family(theme)
        cp.font.color.rgb = title_text_color(theme)

    chart_w = slide_w - 2 * margin
    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("Series", vals)
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        margin,
        chart_top,
        chart_w,
        chart_h,
        chart_data,
    )
    ch = graphic.chart
    ch.has_legend = False
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.font.name = font_family(theme)
    dl.font.size = Pt(10)
    dl.font.color.rgb = theme["NEUTRAL_DARK"]
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    n = len(cats)
    colors = _bar_color_spectrum(theme["SECONDARY_COLOR"], max(1, n))
    for i in range(n):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]

    cat_ax = ch.category_axis
    cat_ax.tick_labels.font.name = font_family(theme)
    cat_ax.tick_labels.font.size = Pt(10)
    cat_ax.tick_labels.font.color.rgb = theme["NEUTRAL_DARK"]
    val_ax = ch.value_axis
    val_ax.has_major_gridlines = True
    val_ax.major_gridlines.format.line.color.rgb = theme["NEUTRAL_LIGHT"]
    val_ax.tick_labels.font.name = font_family(theme)
    val_ax.tick_labels.font.size = Pt(9)

    if sources_line:
        draw_sources_footer(slide, prs, theme, sources_line)
    return slide


if __name__ == "__main__":
    from themes import get_theme

    CITATION_URLS = {
        14: "https://example.com/source/14",
        15: "https://example.com/source/15",
        16: "https://example.com/source/16",
    }

    prs = Presentation()
    ft = get_theme("finance")
    create_finance_dual_column_bar_slide(
        prs,
        title="Secular Tailwinds: Grid Modernization & Data Centers",
        columns=[
            {
                "heading": "Grid Modernization & AI Catalysts",
                "bullets": [
                    {
                        "lead": "Aging Infrastructure",
                        "body": "Transmission replacement cycle accelerating.",
                        "cites": [14],
                    },
                    {
                        "lead": "AI Power Demand",
                        "body": "Data center load growth well above trend.",
                        "cites": [15],
                    },
                ],
            },
            {
                "heading": "Market Opportunity & Structural Shifts",
                "bullets": [
                    {
                        "lead": "Total Addressable Market",
                        "body": "Large specialty rental TAM with runway.",
                        "cites": [16],
                    },
                ],
            },
        ],
        chart={
            "title": "U.S. Data Center Electricity Demand (GW)",
            "categories": ["2025", "2035"],
            "values": [40.0, 106.0],
        },
        sources_line="10-K • Mar 24 [14,15,16]",
        citation_urls=CITATION_URLS,
        theme=ft,
    )

    tests_dir = os.path.join(os.path.dirname(__file__), "..", "tests")
    os.makedirs(tests_dir, exist_ok=True)
    output_path = os.path.join(tests_dir, "test_finance_dual_column_bar_slide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
